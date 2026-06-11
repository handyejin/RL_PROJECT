"""팀 통합 발표자료 초안을 최신 실험 프로토콜 기준으로 정리한다.

기존 PPT의 오래된 마포구/7일/BC 중심 표현을 제거하고,
실험 결과 숫자는 각 담당자가 나중에 넣을 수 있도록 자리만 남긴다.
장표마다 용어 주석과 참고문헌 연결 문장을 넣어 비전공자도 흐름을 따라갈 수 있게 한다.
"""

from __future__ import annotations

from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SOURCE = Path(
    "/Users/jaeyoung/Desktop/RL_PROJECT_2/output/ppt/"
    "ddareungi_team_presentation_no_results_2026-06-11_0847.pptx"
)
STAMP = datetime.now().strftime("%Y-%m-%d_%H%M")
OUT = Path(
    "/Users/jaeyoung/Desktop/RL_PROJECT_2/output/ppt/"
    f"ddareungi_team_presentation_mdp_rich_{STAMP}.pptx"
)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(28, 45, 71)
BLUE = RGBColor(42, 106, 220)
MINT = RGBColor(25, 135, 101)
DARK = RGBColor(31, 41, 55)
GRAY = RGBColor(107, 114, 128)
LIGHT = RGBColor(255, 255, 255)
LINE = RGBColor(219, 226, 239)
RED = RGBColor(210, 74, 67)
PALE_YELLOW = RGBColor(255, 249, 225)


def clear_slide(slide):
    """슬라이드의 기존 도형을 제거한다."""

    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=18,
    bold=False,
    color=DARK,
    align=PP_ALIGN.LEFT,
    font="Apple SD Gothic Neo",
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None, page=None):
    add_text(slide, title, 0.55, 0.35, 9.8, 0.45, size=24, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, subtitle, 0.58, 0.86, 11.5, 0.35, size=11.5, color=GRAY)
    if page:
        add_text(slide, page, 11.9, 0.42, 0.8, 0.25, size=10, color=GRAY, align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(1.18), Inches(12.2), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.color.rgb = LINE


def add_card(slide, x, y, w, h, title, body, *, accent=BLUE):
    rect = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = LIGHT
    rect.line.color.rgb = LINE
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.06), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.color.rgb = accent
    add_text(slide, title, x + 0.18, y + 0.15, w - 0.35, 0.25, size=13.5, bold=True, color=NAVY)
    add_text(slide, body, x + 0.18, y + 0.52, w - 0.35, h - 0.65, size=10.5, color=DARK)


def add_note(slide, text, x=0.72, y=6.72, w=11.9, h=0.42, *, color=GRAY):
    """장표 하단에 발표자가 읽을 수 있는 짧은 메시지나 용어 주석을 넣는다."""

    rect = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = PALE_YELLOW
    rect.line.color.rgb = RGBColor(238, 225, 165)
    add_text(slide, text, x + 0.12, y + 0.08, w - 0.24, h - 0.1, size=8.7, color=DARK)


def add_message(slide, text, *, y=6.34):
    add_text(slide, text, 0.86, y, 11.6, 0.28, size=10.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def add_bullets(slide, items, x, y, w, h, *, size=13, color=DARK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Apple SD Gothic Neo"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return box


def add_table_like(slide, x, y, w, row_h, headers, rows, widths=None, size=9.2):
    if widths is None:
        widths = [w / len(headers)] * len(headers)
    cx = x
    for j, head in enumerate(headers):
        rect = slide.shapes.add_shape(1, Inches(cx), Inches(y), Inches(widths[j]), Inches(row_h))
        rect.fill.solid()
        rect.fill.fore_color.rgb = NAVY
        rect.line.color.rgb = RGBColor(255, 255, 255)
        add_text(slide, head, cx + 0.05, y + 0.07, widths[j] - 0.1, row_h - 0.1, size=size, bold=True, color=RGBColor(255,255,255), align=PP_ALIGN.CENTER)
        cx += widths[j]
    for i, row in enumerate(rows):
        cx = x
        yy = y + row_h * (i + 1)
        for j, val in enumerate(row):
            rect = slide.shapes.add_shape(1, Inches(cx), Inches(yy), Inches(widths[j]), Inches(row_h))
            rect.fill.solid()
            rect.fill.fore_color.rgb = RGBColor(255, 255, 255) if i % 2 == 0 else LIGHT
            rect.line.color.rgb = LINE
            add_text(slide, str(val), cx + 0.06, yy + 0.07, widths[j] - 0.12, row_h - 0.1, size=size, color=DARK)
            cx += widths[j]


def build_deck():
    prs = Presentation(SOURCE)
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    while len(prs.slides) < 18:
        prs.slides.add_slide(prs.slide_layouts[0])

    for slide in prs.slides:
        clear_slide(slide)

    # 1
    s = prs.slides[0]
    add_text(s, "수요예측과 후보 행동 구조를 이용한 서울 따릉이 재배치 강화학습", 0.7, 1.05, 11.8, 0.8, size=28, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, "서울 25개 구 · 73일 chronological holdout · REINFORCE / A2C / DQN 계열 / PPO 비교", 1.3, 2.0, 10.7, 0.45, size=16, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(s, "박제영(A73024), 손예진(A73006), 이형진(A73031)", 2.6, 3.0, 8.2, 0.35, size=14, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, "GitHub: handyejin/RL_PROJECT", 3.4, 3.42, 6.5, 0.3, size=12, color=BLUE, align=PP_ALIGN.CENTER)
    add_card(s, 1.2, 4.45, 10.9, 1.25, "발표 핵심", "정류소 재고 불균형을 줄이는 순차 의사결정 문제로 따릉이 재배치를 정의하고, 수요예측 feature와 Top-K 후보 행동 구조가 강화학습 성능과 안정성에 어떤 영향을 주는지 검증한다.", accent=MINT)
    add_note(s, "발표 흐름: 문제 제기 → MDP 정의 → State/Action 개선 → 실험 설계 → 알고리즘별 결과 삽입 → 활용방안")

    # 2
    s = prs.slides[1]
    add_title(s, "팀 구성 및 역할분담", "각 담당자가 수행한 구현·실험 범위를 최신 기준으로 정리", "02 / 18")
    add_card(s, 0.7, 1.55, 3.85, 4.8, "손예진", "환경/데이터 전처리\nState·Action·Reward 설계\nReplay Viewer 및 시각화\n공통 실행/README 정리", accent=BLUE)
    add_card(s, 4.75, 1.55, 3.85, 4.8, "박제영", "REINFORCE with Value Baseline\nA2C(1-step TD Actor-Critic)\nVAE latent / Contextual Bandit\nTop-K·Seed·Ablation 실험 정리", accent=MINT)
    add_card(s, 8.8, 1.55, 3.85, 4.8, "이형진", "PPO / QRDQN / DQN 계열 실험\nSeed 반복 및 검증 지표\n통합 결과 정리\n최종 발표자료 구성", accent=RED)
    add_note(s, "역할분담은 환경·데이터, 정책 학습, 안정성 검증, 시각화/보고서 산출물 기준으로 정리했다.")

    # 3
    s = prs.slides[2]
    add_title(s, "문제 정의: 왜 강화학습인가?", "따릉이 재배치는 예측 문제를 넘어, 행동이 미래 상태를 바꾸는 순차 의사결정 문제", "03 / 18")
    add_card(s, 0.62, 1.36, 3.95, 4.65, "확보한 데이터와 문제의 성격", "확보 가능 데이터:\n· 정류소별 10분 대여/반납 이력\n· 정류소 위치, capacity, 구 정보\n· 시간/요일/날씨 정보\n· ML 기반 1시간 수요예측\n\n이 데이터로 하루 episode를 재생할 수 있고, 트럭 이동 후 다음 재고 상태를 계산할 수 있다. 그래서 단순 예측이 아니라 시뮬레이터 기반 MDP로 만들 수 있었다.", accent=BLUE)
    add_card(s, 4.72, 1.36, 3.95, 4.65, "왜 수요예측을 먼저 했나", "서울 전체 정류소를 그대로 action으로 두면 후보가 너무 많아진다. RL은 많은 action 중에서 시행착오로 좋은 이동을 찾아야 하므로 탐색이 매우 어려워진다.\n\n그래서 과거 대여/반납·시간·요일·날씨로 1시간 수요를 먼저 예측하고, 이 예측값으로 ‘어디가 곧 부족/과잉이 될지’를 State와 Top-K 후보 생성에 넣었다.", accent=MINT)
    add_card(s, 8.82, 1.36, 3.95, 4.65, "논문 근거", "KDD 2018: hourly station-level rentals/returns 예측이 effective rebalancing의 핵심 입력이라고 설명.\n\nLiang et al.(2024): dynamic rebalancing을 MDP로 공식화하고 reward를 lost demand의 음수로 정의.\n\nSutton & Barto: RL은 trial-and-error search와 delayed reward가 핵심 특징.\n\n따라서 우리는 예측값을 State로 쓰고, 재배치 선택은 RL policy로 학습했다.", accent=RED)
    add_message(s, "핵심 메시지: ML은 서울 전체 수요 압력을 요약하고, RL은 그 정보를 써서 지금 이동할 정류소를 결정한다.", y=6.08)
    add_note(s, "용어: MDP는 State-Action-Reward-Next State로 순차 의사결정을 표현하는 형식이다. Top-K는 전체 정류소 중 후보 K개만 먼저 뽑아 action을 줄이는 구조다.")

    # 4
    s = prs.slides[3]
    add_title(s, "우리가 검증하려는 것", "알고리즘 이름보다 먼저, 어떤 설계가 학습을 가능하게 했는지 확인", "04 / 18")
    add_card(s, 0.85, 1.35, 5.7, 1.28, "질문 1. 무엇을 보게 할 것인가?", "현재 재고만으로 충분한가, 아니면 ML 기반 1시간 수요예측 feature가 필요한가?", accent=BLUE)
    add_card(s, 0.85, 2.78, 5.7, 1.28, "질문 2. 어떤 행동공간이 학습 가능한가?", "서울 전체 정류소를 직접 고르면 action이 너무 커진다. 수요예측 기반 Top-K 후보 rank로 줄이면 학습이 가능한가?", accent=MINT)
    add_card(s, 0.85, 4.21, 5.7, 1.28, "질문 3. 어떤 학습 방식이 안정적인가?", "MC 기반 REINFORCE, TD 기반 A2C, clipped PPO, value-based DQN의 안정성을 비교한다.", accent=RED)
    add_bullets(s, [
        "최종 발표에서는 ‘모든 알고리즘이 항상 이겼다’가 아니라, 어떤 설계에서 어떤 알고리즘이 안정적인지 보여준다.",
        "비교 기준은 73일 holdout의 MostImbalanced 대비 Delta로 통일한다.",
        "초기 탐색 실험은 본문 결론의 근거로 쓰지 않고, 최신 공통 protocol 결과만 중심에 둔다."
    ], 7.0, 1.55, 5.3, 3.7, size=12.6)
    add_message(s, "핵심 차별점: 수요예측은 정답 행동이 아니라, 서울 전체 정류소 중 곧 불균형이 커질 후보를 압축하는 정보다.", y=6.08)
    add_note(s, "용어: baseline은 비교 기준 정책이다. 본 발표의 baseline은 MostImbalanced로, 목표 재고에서 가장 많이 벗어난 정류소를 우선 방문하는 규칙이다.")

    # 5
    s = prs.slides[4]
    add_title(s, "환경 및 데이터셋", "서울 25개 구, 시간순 학습/평가 분할, 공공데이터 기반 전처리", "05 / 18")
    headers = ["구분", "사이트명 / URL", "본 프로젝트에서의 용도"]
    rows = [
        ["대여 이력", "서울 열린데이터광장 OA-15182\nhttps://data.seoul.go.kr/dataList/OA-15182/F/1/datasetView.do", "정류소별 rentals/returns 집계"],
        ["대여소 정보", "서울 열린데이터광장 OA-13252\nhttps://data.seoul.go.kr/dataList/OA-13252/F/1/datasetView.do", "좌표, 관리번호, capacity"],
        ["실시간 재고", "서울 열린데이터광장 OA-15493\nhttps://data.seoul.go.kr/dataList/OA-15493/A/1/datasetView.do", "향후 실제 운영 입력 후보"],
        ["기상 정보", "기상청 ASOS\nhttps://data.kma.go.kr/data/grnd/selectAsosRltmList.do", "날씨 feature / 수요예측 보조"],
        ["수요예측", "HistGradientBoostingRegressor\n입력: station/time/weather/past demand/profile\n출력: rentals/returns/net 1h", "서울 전체 action 후보를 압축하는 State/Top-K feature"],
        ["학습 산출물", "data/processed_seoul_all, data/forecast_by_gu", "episode, forecast, cache"],
    ]
    add_table_like(s, 0.55, 1.27, 12.25, 0.49, headers, rows, widths=[1.55, 6.25, 4.45], size=7.0)
    add_text(s, "수요예측 이유: 전체 정류소 action은 너무 크므로, ML이 1시간 수요 압력을 예측해 RL의 State와 Top-K 후보 생성을 돕는다.", 0.8, 5.72, 11.8, 0.28, size=9.4, bold=True, color=NAVY)
    add_text(s, "분할: 2025년 앞 80% train, 2025-10-20~12-31 73일 holdout eval. Random split은 미래 패턴 누수 위험이 있어 최종 비교에서 제외한다.", 0.8, 6.04, 11.8, 0.32, size=10.0, bold=True, color=NAVY)
    add_note(s, "용어: holdout은 학습에 쓰지 않고 마지막 평가에만 남겨 둔 데이터다. 이번 기준은 2025-10-20~12-31, 73일이다.")

    # 6
    s = prs.slides[5]
    add_title(s, "수요예측 모델링", "ML forecast로 서울 전체 수요 압력을 요약하고, RL의 State와 Top-K 후보 생성에 연결", "06 / 18")
    add_card(s, 0.6, 1.35, 3.95, 4.85, "입력 데이터 예시", "10분 단위 정류소별 수요 집계\n\nCSV/parquet 형태 예시:\n\n t              station  rentals returns\n 2025-01-01    ST-1017       1       0\n 2025-01-01    ST-1035       1       0\n 2025-01-01    ST-1059       1       0\n\n추가 feature:\nis_weekend, holiday, 시간/요일, 과거 1h·3h 수요, 정류소별 profile", accent=BLUE)
    add_card(s, 4.7, 1.35, 3.95, 4.85, "예측 모델", "모델: HistGradientBoostingRegressor\n\nrentals와 returns를 별도 회귀 모델로 학습\n\n입력 X:\nstation_idx, dow, hour, slot, month,\npast_rent_1h, past_ret_1h,\npast_rent_3h, past_ret_3h,\nprofile_rent_1h, profile_ret_1h\n\n목표 y:\n앞으로 6개 10분 구간, 즉 1시간 rentals/returns 합계", accent=MINT)
    add_card(s, 8.8, 1.35, 3.95, 4.85, "예측 결과와 RL 연결", "예측 산출물 예시:\n\n station  pred_rent pred_return pred_net\n ST-1171      0.276      0.187   -0.089\n ST-1172      0.147      0.183    0.036\n ST-1173      0.134      0.234    0.100\n\nRL 사용처:\n1. State에 미래 수요 압력 추가\n2. forecast_imbalance로 Top-K 후보 생성\n3. action space를 전체 정류소에서 K개 rank로 축소", accent=RED)
    add_message(s, "핵심 메시지: 수요예측은 정답 action이 아니라, 넓은 서울 정류소 후보를 RL이 탐색 가능한 크기로 줄이는 정보 압축 단계다.", y=6.16)
    add_note(s, "예측 대상은 실제 미래를 직접 보는 oracle이 아니라, train 기간의 과거 수요·시간 패턴으로 학습한 1시간 ahead forecast다.")

    # 7
    s = prs.slides[6]
    add_title(s, "Hybrid 구조: 지도학습 예측 + 강화학습 정책", "수요예측 모델과 RL agent를 분리해, 예측은 정보로 쓰고 행동은 reward로 학습", "07 / 18")
    add_card(s, 0.62, 1.35, 3.78, 4.82, "1단계: Supervised Forecasting", "목표:\n정류소별 향후 1시간 rentals/returns 예측\n\n입력:\n과거 대여/반납, 시간/요일, 정류소 profile, 날씨\n\n모델:\nHistGradientBoostingRegressor\n\n산출:\npred_rentals_1h, pred_returns_1h, pred_net_1h", accent=BLUE)
    add_card(s, 4.78, 1.35, 3.78, 4.82, "2단계: RL Policy Learning", "목표:\n누적 reward가 좋은 재배치 정책 학습\n\n입력 State:\n현재 재고 + capacity + 트럭 상태 + 수요예측 feature\n\nAction:\nTop-K 후보 rank 선택\n\nReward:\nstockout/full/travel cost 최소화", accent=MINT)
    add_card(s, 8.95, 1.35, 3.78, 4.82, "근거와 해석", "KDD 2018은 정류소별 hourly rentals/returns 예측이 effective rebalancing의 핵심 입력이라고 설명한다.\n\n최근 PPO 기반 bike rebalancing 연구도 historical usage, station attributes, weather, demand forecasts를 state로 사용한다.\n\n즉 예측 모델은 정책을 대체하지 않고, RL이 볼 수 있는 미래 수요 정보를 제공한다.", accent=RED)
    add_message(s, "핵심 메시지: 이 프로젝트는 ‘지도학습 vs 강화학습’이 아니라, 지도학습 예측을 RL 의사결정에 연결한 Hybrid pipeline이다.", y=6.14)
    add_note(s, "용어: Hybrid는 수요예측 모델이 행동을 직접 결정하지 않고, 예측값을 RL의 State/후보 생성 feature로 넘겨 정책 학습을 돕는 구조를 뜻한다.")

    # 8
    s = prs.slides[7]
    add_title(s, "MDP 설계: State · Action · Reward", "과제 평가의 핵심: 문제를 어떤 MDP로 만들었는지 명확히 설명", "08 / 18")
    add_card(s, 0.6, 1.38, 3.95, 4.7, "State: 무엇을 보나", "현재 재고 비율 bikes/capacity\n정류소별 capacity\n트럭 위치·적재량\n시간/요일/날씨\nML 예측값: pred_rentals_1h, pred_returns_1h, pred_net_1h\n\n구별 정류소 수가 달라 obs_dim도 다르며, 각 구별 독립 episode로 학습/평가한다.", accent=BLUE)
    add_card(s, 4.7, 1.38, 3.95, 4.7, "Action: 무엇을 선택하나", "서울 전체 정류소 직접 선택 대신 Top-K 후보 rank 선택\n\ncandidate_score = forecast_imbalance\n  - 0.20 × distance\n  - 1.0 × zone_penalty\n\nDulac-Arnold et al.(2015)는 큰 이산 action space가 기존 RL 적용을 어렵게 만든다고 지적한다. Action Elimination 연구도 불필요한 action 제거가 학습을 개선한다고 보고한다.", accent=MINT)
    add_card(s, 8.8, 1.38, 3.95, 4.7, "Reward: 무엇을 줄이나", "r_t = -1.0×stockout\n      -0.8×full\n      -0.008×travel_km\n      -0.002×travel_step\n\n평가 시 urgent_bonus=0, shaping_scale=0으로 공정 metric만 사용한다.\n\n근거: 동적 재배치 연구의 목표는 lost demand / availability failure를 줄이는 것이다.", accent=RED)
    add_message(s, "참고문헌 연결: State는 수요예측 연구, Action은 large discrete action space 연구, Reward는 lost demand 최소화 연구와 직접 연결된다.", y=6.17)
    add_note(s, "용어: Top-K는 정답을 주는 알고리즘이 아니라, 전체 정류소 중 후보 K개를 뽑아 RL이 rank를 선택하게 하는 action 구조다.")

    # 9
    s = prs.slides[8]
    add_title(s, "Reward가 발생하는 상황 예시", "좋은 행동은 대여/반납 실패와 불필요한 이동을 줄여 누적 reward를 높인다", "09 / 18")
    add_card(s, 0.75, 1.42, 3.75, 4.35, "Stockout penalty", "정류소에 자전거가 부족해 대여 수요를 처리하지 못하면 음수 보상.\n\n예: 출근 시간 수요가 큰 정류소를 방치하면 다음 step에서 stockout 증가.\n\n논문 연결: lost rentals / unsatisfied demand를 줄이는 것이 BSS 운영자의 핵심 목표.", accent=RED)
    add_card(s, 4.8, 1.42, 3.75, 4.35, "Full penalty", "거치대가 가득 차 반납 수요를 처리하지 못하면 음수 보상.\n\n예: 반납이 몰리는 역세권을 비워두지 않으면 full 증가.\n\n논문 연결: full or empty stations가 서비스 신뢰도를 떨어뜨리는 주요 원인.", accent=BLUE)
    add_card(s, 8.85, 1.42, 3.75, 4.35, "Travel cost", "먼 정류소로 이동할수록 비용이 증가.\n\n모델은 큰 불균형을 줄이되 이동거리도 함께 고려해야 한다.\n\n논문 연결: vehicle-based rebalancing은 운영비와 routing cost를 함께 고려해야 하는 문제.", accent=MINT)
    add_message(s, "핵심 메시지: reward가 즉시 좋은/나쁜 행동을 알려주지 않기 때문에 credit assignment가 어렵다.", y=6.05)
    add_note(s, "용어: episode는 하루 시뮬레이션 한 번이다. step마다 대여/반납 수요와 트럭 이동을 반영해 next state를 계산하고, 하루 reward 합으로 평가한다.")

    # 10
    s = prs.slides[9]
    add_title(s, "비교 알고리즘 개요", "각 알고리즘의 업데이트 방식과 기대 특성을 분리해서 설명", "10 / 18")
    headers = ["알고리즘", "핵심 방식", "업데이트 단위", "발표 포인트"]
    rows = [
        ["REINFORCE", "Reward-to-go + Value baseline", "Episode 종료 후", "MC라 seed/trajectory 분산이 큼"],
        ["A2C", "Actor-Critic, 1-step TD", "Step/batch", "Actor/Critic 독립 optimizer"],
        ["DQN 계열", "Double DQN / QRDQN", "Replay sample", "rank action에서 Q target 학습 난이도"],
        ["PPO", "Clipped actor-critic", "Rollout batch", "policy update 변화 폭 제한"],
        ["LinUCB", "Contextual Bandit", "즉시 reward", "장기 return을 보지 못하는 한계"],
    ]
    add_table_like(s, 0.55, 1.34, 12.2, 0.5, headers, rows, widths=[2.0, 3.5, 2.4, 4.3], size=8.2)
    add_message(s, "참고문헌 연결: PPO는 clipped surrogate objective로 policy update가 과도하게 바뀌지 않도록 제한한다.", y=5.98)
    add_note(s, "용어: Seed는 난수 초기값이다. 같은 코드라도 seed가 달라지면 초기 가중치와 sampling 순서가 달라져 결과가 흔들릴 수 있다.")

    # 11
    s = prs.slides[10]
    add_title(s, "실험 설계: Sequential Screening", "Train reward가 아니라 holdout 평가 reward로 후보를 좁히고 확인", "11 / 18")
    steps = [
        ("1. Full baseline", "25개 구 전체, 기본 Top-K, seed 42로 전체 경향 확인"),
        ("2. Best/Worst 선정", "Best Delta 기준 상·하위 구를 뽑아 민감도 분석 대상 구성"),
        ("3. Top-K ablation", "K=3/6/9/12/15를 비교해 action 후보 수 영향 확인"),
        ("4. Confirmation", "선택된 K를 500 episode 또는 충분한 step으로 재확인"),
        ("5. Seed validation", "seed 42/123/777 반복으로 평균·표준편차 확인"),
        ("6. Final run", "선택 조건으로 전체 25개 구 최종 실행"),
    ]
    for i, (t, b) in enumerate(steps):
        x = 0.75 + (i % 3) * 4.05
        y = 1.38 + (i // 3) * 2.22
        add_card(s, x, y, 3.65, 1.62, t, b, accent=[BLUE, MINT, RED][i % 3])
    add_message(s, "핵심 메시지: 학습 중 train reward가 높아도 일반화가 보장되지 않으므로, 별도 holdout episode 평가로 비교한다.", y=5.95)
    add_note(s, "용어: Train reward는 학습에 사용한 episode의 보상이고, holdout eval reward는 학습에 쓰지 않은 날짜에서 정책을 실행해 얻은 보상이다.")

    # 12
    s = prs.slides[11]
    add_title(s, "평가 지표", "Train reward가 아니라 holdout episode 평가 reward로 비교", "12 / 18")
    add_table_like(s, 0.85, 1.35, 11.6, 0.48, ["지표", "정의", "왜 필요한가"], [
        ["Mean Reward", "73일 holdout 평균 누적 reward", "평가 기간 전체 성능"],
        ["Delta", "Model reward - MostImbalanced reward", "구별 reward scale 차이 완화"],
        ["Best checkpoint", "중간 평가 중 최고 모델", "논문식 model selection"],
        ["Final checkpoint", "마지막 모델", "학습 안정성 확인"],
        ["Best-Final gap", "Best와 Final 차이", "forgetting/불안정성 진단"],
        ["Seed std", "같은 구에서 seed별 표준편차", "재현성과 민감도 평가"],
        ["PPO diagnostics", "approx_kl, clip_fraction, entropy", "clipped update와 정책 변화량 확인"],
    ], widths=[2.25, 4.65, 4.7], size=8.0)
    add_note(s, "용어: 승리 구 수는 Delta > 0인 구의 개수다. Best checkpoint는 학습 중 평가가 가장 좋았던 모델, Final checkpoint는 마지막 모델이다.")

    # 13
    s = prs.slides[12]
    add_title(s, "전체 결과 표 삽입 위치", "이 장은 각 담당자가 최신 73일 holdout 결과를 넣는 공통 템플릿", "13 / 18")
    add_card(s, 0.8, 1.38, 5.8, 4.5, "넣을 표", "알고리즘별 행:\nREINFORCE / A2C / DQN 계열 / PPO / Bandit\n\n컬럼:\nBest Reward, Final Reward,\nMostImbalanced Reward,\nBest Delta, Final Delta,\n승리 구 수(Delta>0), Best step/episode", accent=BLUE)
    add_card(s, 6.9, 1.38, 5.6, 4.5, "주의", "실험결과 숫자는 여기서 고정하지 않는다.\n\n각 담당자가 동일 기준으로 산출한 최신 값을 삽입한다.\n\n기준은 반드시 73일 chronological holdout, Delta = model - baseline으로 통일한다.", accent=RED)
    add_note(s, "용어: MostImbalanced는 현재 재고가 목표 재고에서 가장 크게 벗어난 정류소를 우선 방문하는 규칙 기반 baseline이다.")

    # 14
    s = prs.slides[13]
    add_title(s, "REINFORCE / A2C 결과 해석 장표", "박제영 담당 결과를 넣을 때의 권장 구성", "14 / 18")
    add_card(s, 0.75, 1.35, 3.8, 4.55, "REINFORCE", "Policy gradient를 episode 단위로 업데이트.\n\nReward-to-go와 Value baseline으로 variance를 낮췄지만, trajectory seed에 민감할 수 있음.\n\n넣을 그림: seed별 Delta scatter 또는 error bar.", accent=BLUE)
    add_card(s, 4.78, 1.35, 3.8, 4.55, "A2C", "Actor와 Critic을 함께 학습.\n\n1-step TD target으로 advantage를 추정하고, actor_loss와 critic_loss는 독립 optimizer로 각각 업데이트한다.\n\n넣을 그림: A2C vs REINFORCE seed std 비교.", accent=MINT)
    add_card(s, 8.8, 1.35, 3.8, 4.55, "해석 문장", "A2C가 항상 최고라는 주장보다, TD 기반 업데이트가 이 환경에서 seed 민감도를 낮추는 경향을 보였는지 중심으로 설명한다.", accent=RED)
    add_note(s, "용어: MC는 episode가 끝난 뒤 실제 return으로 업데이트한다. TD는 다음 value 예측을 이용해 더 자주 업데이트한다.")

    # 15
    s = prs.slides[14]
    add_title(s, "PPO / DQN 결과 해석 장표", "팀원 담당 결과를 같은 평가 기준으로 통합", "15 / 18")
    add_card(s, 0.8, 1.35, 5.7, 4.75, "PPO", "핵심은 clipped surrogate objective.\n\n넣을 지표:\napprox_kl, clip_fraction, entropy, explained_variance\n\n해석:\nold policy와 new policy의 차이가 작게 유지되는지 확인한다. clip_fraction은 clipping이 적용된 sample 비율이므로 낮을수록 clipping 발동이 적고, 높을수록 update가 더 자주 제한된 것으로 해석한다.", accent=BLUE)
    add_card(s, 6.85, 1.35, 5.7, 4.75, "DQN / QRDQN", "핵심은 value-based off-policy 학습.\n\n넣을 지표:\nBest/Final Delta, replay 기반 학습곡선, target Q 안정성, Top-K별 차이\n\n해석:\nrank action에서 Q target 학습이 어려웠는지 확인.", accent=MINT)
    add_note(s, "용어: PPO clipping은 old policy와 new policy의 action 확률 비율을 제한해 너무 큰 policy 변화를 막는 장치다.")

    # 16
    s = prs.slides[15]
    add_title(s, "Replay Viewer 시연", "정책이 episode 안에서 어떤 정류소를 선택하고 트럭이 어떻게 이동하는지 보여주는 장표", "16 / 18")
    add_card(s, 0.85, 1.5, 7.0, 4.8, "Replay Viewer에서 확인할 것", "10~20초 시연 권장\n\n보여줄 것:\n1. 정류소 점 색상/크기\n2. Top-K 후보 변화\n3. 선택 정류소와 트럭 이동\n4. 누적 reward 변화\n\n숫자 결과와 실제 정책 행동을 함께 보여준다.", accent=BLUE)
    add_card(s, 8.2, 1.5, 4.25, 4.8, "설명 포인트", "그래프 숫자만으로 보이지 않는 정책 행동을 시각적으로 검증한다.\n\n특정 정류소를 계속 고르는 collapse나 후보 분포 변화를 확인할 수 있다.", accent=MINT)
    add_note(s, "발표 팁: 영상은 ‘예쁜 시각화’보다 정책이 실제로 어떤 후보를 보고 어떤 정류소로 이동하는지 설명하는 증거로 사용한다.")

    # 17
    s = prs.slides[16]
    add_title(s, "활용방안", "따릉이를 넘어, 공공 수요예측 + 순차 의사결정 문제로 확장 가능", "17 / 18")
    add_card(s, 0.75, 1.35, 3.8, 4.9, "따릉이 운영 보조", "다음 1시간 수요를 반영해 어떤 정류소를 우선 방문할지 후보를 제안한다.\n\n사람 운영자가 최종 판단하는 재배치 우선순위 추천 도구로 활용 가능.", accent=BLUE)
    add_card(s, 4.78, 1.35, 3.8, 4.9, "공공데이터 확장 사례", "공영주차장 혼잡 완화\n전기차 충전소 대기 관리\n공공자전거/킥보드 재배치\n버스·택시 수요 대응\n\n공통점: 수요가 시간에 따라 변하고, 현재 배치가 미래 비용을 바꿈.", accent=MINT)
    add_card(s, 8.8, 1.35, 3.8, 4.9, "실서비스 전 조건", "실시간 재고 API\n차량 수/기사 근무시간\n도로 이동시간\n안전·교통 규칙\n정책 설명 가능성\n\n현재는 자동 운영이 아니라 simulator 기반 의사결정 보조.", accent=RED)
    add_note(s, "활용 범위: 공유자원 배치, 혼잡 완화, 재고·인력 재배치처럼 ‘예측 후 순차적으로 움직여야 하는’ 공공서비스 문제에 응용할 수 있다.")

    # 18
    s = prs.slides[17]
    add_title(s, "결론 및 참고문헌", "결과 수치보다 설계와 검증 기준을 명확히 정리", "18 / 18")
    add_card(s, 0.75, 1.3, 5.6, 4.95, "결론 문장", "따릉이 재배치는 정류소 재고와 미래 수요가 함께 변하는 순차 의사결정 문제다.\n\nML 수요예측은 미래 대여/반납 압력을 State로 제공하고, Top-K 후보 행동 구조는 RL이 탐색 가능한 문제 크기로 줄이는 핵심 설계였다.\n\n결과 해석은 73일 holdout의 Delta, Best-Final gap, seed std를 함께 보고 알고리즘별 안정성 차이로 정리한다.", accent=BLUE)
    refs = [
        "Liang et al. (2024): BSS는 stochastic demand로 full/empty station이 생기며, dynamic rebalancing을 MDP로 정의 → 3p/6p 근거",
        "KDD 2018 demand prediction: hourly station-level rentals/returns를 temporal/weather feature로 예측 → 5p/6p State 근거",
        "PPO bike rebalancing / data-driven rebalancing studies: historical usage, station attributes, weather, demand forecasts를 의사결정 입력으로 사용 → 7p Hybrid 근거",
        "Dulac-Arnold et al. (2015), Action Elimination (2018): large discrete action space는 탐색/계산을 어렵게 하며, 후보 축소가 학습 효율을 높일 수 있음 → 4p/6p Top-K 근거",
        "Schulman et al. (2017): PPO clipped surrogate objective → 8p/13p PPO 안정화 설명 근거",
        "Henderson et al. (2018): RL seed variance 보고 → 4p/10p seed 평가 근거",
        "Sutton & Barto (2018): MDP, return, TD/MC 기본 개념 → 3p/8p 용어 근거",
    ]
    add_bullets(s, refs, 6.75, 1.32, 5.8, 3.35, size=8.4, color=DARK)
    add_text(s, "데이터: 서울 열린데이터광장 대여이력/대여소 정보/실시간 대여정보, 기상청 ASOS", 6.75, 4.95, 5.8, 0.5, size=8.9, color=GRAY)
    add_text(s, "참고문헌은 ‘권위 있어 보이기’가 아니라, 왜 State·Action·Reward와 평가 protocol을 이렇게 설계했는지 설명하기 위해 사용했다.", 6.75, 5.55, 5.8, 0.55, size=8.8, color=NAVY)
    add_note(s, "마지막 메시지: ‘알고리즘만 바꾼 실험’이 아니라, 따릉이 문제를 RL이 풀 수 있는 MDP와 평가 프로토콜로 만든 것이 핵심이다.")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    out = build_deck()
    print(out)
