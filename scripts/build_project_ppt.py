"""따릉이 강화학습 프로젝트 발표용 PPT 초안을 생성한다.

보고서의 긴 설명을 발표 슬라이드 구조로 압축한다. 흰 배경, 얇은 선,
큰 제목을 기본 스타일로 사용하고, State/Action/Reward와 실험 결과가
쉽게 보이도록 표와 그림 중심으로 구성한다.
"""

from __future__ import annotations

from datetime import datetime
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "output" / "ppt"
FIG_DIR = ROOT / "docs" / "figures"

FONT = "Apple SD Gothic Neo"
TITLE = RGBColor(28, 35, 50)
BODY = RGBColor(55, 65, 81)
MUTED = RGBColor(107, 114, 128)
BLUE = RGBColor(37, 99, 235)
GREEN = RGBColor(5, 150, 105)
RED = RGBColor(220, 38, 38)
AMBER = RGBColor(217, 119, 6)
LINE = RGBColor(229, 231, 235)
FILL = RGBColor(248, 250, 252)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = BODY,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    """텍스트 상자를 추가한다."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.margin_left = Inches(0.04)
    box.text_frame.margin_right = Inches(0.04)
    box.text_frame.margin_top = Inches(0.03)
    box.text_frame.margin_bottom = Inches(0.03)
    box.text_frame.vertical_anchor = valign
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title: str, subtitle: str | None = None):
    """슬라이드 상단 제목 영역을 만든다."""
    add_text(slide, title, 0.55, 0.32, 12.2, 0.48, size=24, bold=True, color=TITLE)
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(0.9), Inches(12.25), Inches(0.01)
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = LINE
    slide.shapes[-1].line.color.rgb = LINE
    if subtitle:
        add_text(slide, subtitle, 0.6, 0.98, 12.0, 0.34, size=10, color=MUTED)


def add_bullets(slide, bullets: list[str], x: float, y: float, w: float, h: float, *, size: int = 16):
    """간결한 bullet 목록을 추가한다."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.space_after = Pt(7)
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = BODY
    return box


def add_card(slide, title: str, body: str, x: float, y: float, w: float, h: float, *, accent=BLUE):
    """정보 카드 하나를 추가한다."""
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = FILL
    shape.line.color.rgb = LINE
    shape.adjustments[0] = 0.08
    add_text(slide, title, x + 0.16, y + 0.13, w - 0.32, 0.28, size=13, bold=True, color=accent)
    add_text(slide, body, x + 0.16, y + 0.48, w - 0.32, h - 0.62, size=11, color=BODY)
    return shape


def add_table(slide, df: pd.DataFrame, x: float, y: float, w: float, h: float, *, font_size: int = 9):
    """DataFrame을 PowerPoint 표로 추가한다."""
    rows, cols = df.shape[0] + 1, df.shape[1]
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for col in range(cols):
        table.cell(0, col).text = str(df.columns[col])
        table.cell(0, col).fill.solid()
        table.cell(0, col).fill.fore_color.rgb = RGBColor(239, 246, 255)
    for r in range(df.shape[0]):
        for c in range(cols):
            table.cell(r + 1, c).text = str(df.iloc[r, c])
    for row in table.rows:
        for cell in row.cells:
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = FONT
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = BODY
    return table


def add_image(slide, path: Path, x: float, y: float, w: float, h: float):
    """이미지를 지정 영역 안에 비율 유지로 배치한다."""
    if not path.exists():
        add_card(slide, "이미지 없음", str(path), x, y, w, h, accent=RED)
        return None
    from PIL import Image

    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    draw_w, draw_h = iw * scale, ih * scale
    px = x + (w - draw_w) / 2
    py = y + (h - draw_h) / 2
    return slide.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(draw_w), Inches(draw_h))


def fmt_num(v, digits=1):
    """표 숫자를 발표용으로 짧게 포맷한다."""
    try:
        return f"{float(v):.{digits}f}"
    except Exception:
        return str(v)


def history_best_final(path: str | Path, metric: str = "eval_reward") -> tuple[float | None, float | None, float | None]:
    """history.npy에서 best/final reward와 best point를 읽는다."""
    import numpy as np

    path = Path(path)
    if not path.exists():
        return None, None, None
    arr = np.load(path, allow_pickle=True)
    rows = [dict(x) for x in arr.tolist()]
    if not rows:
        return None, None, None
    values = [float(row[metric]) for row in rows]
    best_idx = max(range(len(values)), key=lambda i: values[i])
    point_key = "episode" if "episode" in rows[best_idx] else "timesteps"
    return values[best_idx], values[-1], float(rows[best_idx].get(point_key, best_idx))


def new_deck() -> Presentation:
    """프레젠테이션 기본 설정을 만든다."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def blank(prs: Presentation):
    """빈 슬라이드를 추가한다."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def main() -> None:
    algo = pd.read_csv(ROOT / "docs" / "rl_current_algorithm_summary.csv")
    gu = pd.read_csv(ROOT / "docs" / "rl_current_gu_algorithm_summary.csv")
    vae = pd.read_csv(ROOT / "docs" / "rl_current_vae_reinforce_summary.csv")

    prs = new_deck()

    # 1. Title
    slide = blank(prs)
    add_text(slide, "수요예측 기반 따릉이 재배치 강화학습", 0.7, 1.7, 11.8, 0.7, size=30, bold=True, color=TITLE, align=PP_ALIGN.CENTER)
    add_text(slide, "REINFORCE · A2C · PPO · Double DQN · LinUCB Bandit 비교", 1.5, 2.55, 10.3, 0.45, size=16, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "GitHub: handyejin/RL_PROJECT", 4.05, 5.9, 5.2, 0.35, size=12, color=MUTED, align=PP_ALIGN.CENTER)

    # 2. Problem
    slide = blank(prs)
    add_title(slide, "1. 문제 정의", "정류소별 자전거 부족과 반납 불가를 줄이는 순차 의사결정 문제")
    add_card(slide, "현상", "출근·퇴근·주말 수요가 정류소별로 다르게 몰리면서 일부 정류소는 자전거가 부족하고, 일부는 거치 공간이 부족해진다.", 0.75, 1.55, 3.7, 2.0, accent=RED)
    add_card(slide, "의사결정", "재배치 트럭은 현재 위치와 적재량을 가진 상태에서 다음에 방문할 정류소를 반복적으로 선택한다.", 4.8, 1.55, 3.7, 2.0, accent=BLUE)
    add_card(slide, "목표", "하루 episode 동안 stockout, full, 이동 비용을 줄여 누적 reward를 0에 가깝게 만든다.", 8.85, 1.55, 3.7, 2.0, accent=GREEN)
    add_bullets(slide, ["강화학습 관점: 현재 상태를 보고 다음 방문 정류소를 선택하는 MDP", "운영 관점: 대여 실패와 반납 실패를 줄이는 재고 재배치 문제"], 1.0, 4.25, 11.2, 1.2, size=17)

    # 3. Objective
    slide = blank(prs)
    add_title(slide, "2. 최적화 목표", "reward를 최대화한다는 것은 서비스 실패와 이동 비용을 최소화한다는 뜻")
    add_card(slide, "최대화", "7개 평가일 평균 episode reward\nDelta = 모델 reward - MostImbalanced reward", 0.9, 1.55, 5.4, 2.0, accent=GREEN)
    add_card(slide, "최소화", "stockout: 빌릴 자전거가 없음\nfull: 반납할 거치 공간이 없음\ntravel: 불필요한 트럭 이동", 7.0, 1.55, 5.4, 2.0, accent=RED)
    add_text(slide, "rₜ = -1.0·stockout -0.8·full -0.008·travel_km -0.002·travel_step", 1.05, 4.35, 11.25, 0.55, size=22, bold=True, color=TITLE, align=PP_ALIGN.CENTER)
    add_text(slide, "reward는 음수이므로 0에 가까울수록 좋다", 4.15, 5.05, 5.0, 0.32, size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # 4. Data
    slide = blank(prs)
    add_title(slide, "3. 데이터셋과 전처리", "10분 단위 대여/반납 기록을 구별 episode replay 데이터로 변환")
    data_df = pd.DataFrame(
        [
            ["분석 범위", "서울 25개 구"],
            ["정류소", "3,313개"],
            ["Active 정류소", "2,808개"],
            ["대여/반납 row", "40,565,021"],
            ["예측 파일", "구별 25개 parquet"],
        ],
        columns=["항목", "값"],
    )
    add_table(slide, data_df, 0.85, 1.55, 4.3, 2.7, font_size=13)
    add_image(slide, FIG_DIR / "demand_dow_hour.png", 5.55, 1.35, 3.25, 2.7)
    add_image(slide, FIG_DIR / "demand_top_stations.png", 9.1, 1.35, 3.25, 2.7)
    add_bullets(slide, ["원본 기록에서 정류소별 시간 흐름을 만들고, 하루를 하나의 episode로 평가", "실시간 재고 스냅샷이 없으므로 초기 재고에서 대여/반납 event를 반영해 재고를 갱신"], 0.9, 4.85, 11.6, 1.0, size=14)

    # 5. Environment
    slide = blank(prs)
    add_title(slide, "4. 환경 구현", "하루 운영을 10분 단위로 replay하며 트럭이 순차적으로 정류소를 선택")
    add_card(slide, "Episode", "평가 날짜 1일\n10분 단위 수요 replay", 0.8, 1.5, 3.0, 1.6, accent=BLUE)
    add_card(slide, "Inventory", "capacity와 target 대비\n재고 부족/과포화 계산", 4.0, 1.5, 3.0, 1.6, accent=GREEN)
    add_card(slide, "Truck", "현재 위치·적재량·이동거리\n다음 정류소 선택", 7.2, 1.5, 3.0, 1.6, accent=AMBER)
    add_card(slide, "Evaluation", "고정 7개 날짜 평균\nMostImbalanced 대비 Delta", 10.4, 1.5, 2.4, 1.6, accent=RED)
    add_text(slide, "Replay Viewer로 학습된 policy가 실제 episode에서 어떤 정류소를 선택하는지 확인 가능", 1.2, 4.4, 10.9, 0.5, size=19, bold=True, color=TITLE, align=PP_ALIGN.CENTER)
    add_card(slide, "Replay Viewer 활용", "학습된 정책의 정류소 선택, 트럭 이동, reward 변화를 영상으로 확인한다.", 3.2, 5.25, 6.9, 1.0, accent=BLUE)

    # 6. SAR
    slide = blank(prs)
    add_title(slide, "5. State / Action / Reward 설계", "강화학습 문제의 핵심 정의")
    sar = pd.DataFrame(
        [
            ["State", "정류소 재고, capacity, target 편차, 시간, 트럭 위치/적재량, 1시간 수요예측"],
            ["Action", "다음 방문 정류소 선택. 개선 실험에서는 Top-K 후보 12개 중 하나 선택"],
            ["Reward", "stockout/full/이동 비용의 음수 합. 0에 가까울수록 좋음"],
        ],
        columns=["구성", "정의"],
    )
    add_table(slide, sar, 0.8, 1.5, 11.8, 2.5, font_size=12)
    add_text(slide, "projected_bikes = current_bikes + pred_returns_1h - pred_rentals_1h", 1.1, 4.55, 11.1, 0.45, size=18, bold=True, color=TITLE, align=PP_ALIGN.CENTER)
    add_text(slide, "미래 수요를 상태에 넣어, 단순 현재 재고가 아니라 1시간 뒤 부족/과포화 가능성을 보게 한다.", 1.45, 5.2, 10.4, 0.5, size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # 7. Baseline
    slide = blank(prs)
    add_title(slide, "6. Baseline: MostImbalanced", "학습 없이 현재 가장 불균형한 정류소를 방문하는 규칙 기반 정책")
    add_card(slide, "왜 필요한가", "RL 성능을 단순 reward가 아니라 강한 규칙 기반 정책 대비 개선량으로 해석하기 위해 사용한다.", 0.85, 1.55, 5.55, 2.0, accent=BLUE)
    add_card(slide, "평가 방식", "Delta > 0이면 baseline보다 좋고, Delta < 0이면 baseline보다 나쁘다.", 6.95, 1.55, 5.55, 2.0, accent=GREEN)
    add_text(slide, "Delta = model_eval_reward - MostImbalanced_eval_reward", 1.35, 4.5, 10.6, 0.5, size=23, bold=True, color=TITLE, align=PP_ALIGN.CENTER)
    add_text(slide, "구별 reward scale이 다르므로 raw reward보다 Delta와 Delta%를 함께 본다.", 2.55, 5.2, 8.3, 0.35, size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # 8. Improvements
    slide = blank(prs)
    add_title(slide, "7. 문제 발견과 개선 방향", "기본 RL이 어려웠던 이유를 환경 입력과 행동 구조에서 보완")
    add_card(slide, "문제 1: 미래 수요 부족", "현재 재고만 보면 곧 부족해질 정류소를 선제적으로 보기 어렵다.", 0.75, 1.45, 3.7, 1.65, accent=RED)
    add_card(slide, "개선 1: 수요예측 State", "구별 1시간 대여/반납 예측을 observation에 추가했다.", 0.75, 3.35, 3.7, 1.65, accent=GREEN)
    add_card(slide, "문제 2: 행동 공간 큼", "전체 정류소를 직접 선택하면 탐색이 느리고 reward 신호가 흐려진다.", 4.85, 1.45, 3.7, 1.65, accent=RED)
    add_card(slide, "개선 2: Top-K Action", "예측 불균형이 큰 후보 12개 안에서 policy가 선택하도록 했다.", 4.85, 3.35, 3.7, 1.65, accent=GREEN)
    add_card(slide, "문제 3: delayed reward", "한 번의 이동 효과가 뒤늦게 나타나 credit assignment가 어렵다.", 8.95, 1.45, 3.7, 1.65, accent=RED)
    add_card(slide, "개선 3: 평가 분리", "Best와 Final을 분리해 가능성과 안정성을 함께 분석했다.", 8.95, 3.35, 3.7, 1.65, accent=GREEN)

    # 9. Algorithm overview
    slide = blank(prs)
    add_title(slide, "8. 비교한 알고리즘", "정책 경사, actor-critic, value-based, LinUCB 계열을 같은 protocol에서 비교")
    alg_df = pd.DataFrame(
        [
            ["REINFORCE", "On-policy", "episode 후", "policy + value baseline"],
            ["A2C", "On-policy", "step/batch", "actor + critic"],
            ["PPO", "On-policy", "rollout", "clipped policy update"],
            ["Double DQN", "Off-policy", "replay buffer", "Q-network + target"],
            ["LinUCB Bandit", "Bandit", "즉시 업데이트", "후보 feature + uncertainty"],
        ],
        columns=["알고리즘", "종류", "업데이트", "핵심"],
    )
    add_table(slide, alg_df, 0.55, 1.35, 12.25, 4.2, font_size=11)
    add_text(slide, "DQN은 Double DQN과 Dueling Q-network를 사용했고, PPO는 MaskablePPO로 action mask를 적용했다.", 0.95, 6.0, 11.5, 0.4, size=13, color=MUTED, align=PP_ALIGN.CENTER)

    # 10. Loss
    slide = blank(prs)
    add_title(slide, "9. 주요 Loss 설계", "보고서에는 핵심 수식만 짧게 제시")
    add_card(slide, "REINFORCE", "Gₜ = reward-to-go\nAₜ = Gₜ - V(sₜ)\nLoss = -logπ(aₜ|sₜ) Aₜ", 0.75, 1.45, 3.0, 2.0, accent=BLUE)
    add_card(slide, "A2C", "target = r + γV(s')\nA = target - V(s)\nactor + critic loss", 3.95, 1.45, 3.0, 2.0, accent=GREEN)
    add_card(slide, "PPO", "ratio = πnew / πold\nclip(ratio, 1-ε, 1+ε)\n큰 policy 변화 제한", 7.15, 1.45, 3.0, 2.0, accent=AMBER)
    add_card(slide, "Double DQN", "a* = argmax Qonline(s',a)\ny = r + γQtarget(s',a*)\nHuber TD loss", 10.35, 1.45, 2.25, 2.0, accent=RED)
    add_card(slide, "LinUCB", "score = θᵀx + α·uncertainty\n장기 return 없이 즉시 reward만 학습", 3.0, 4.15, 7.4, 1.35, accent=BLUE)

    # 11. Setup
    slide = blank(prs)
    add_title(slide, "10. 실험 셋업", "서울 25개 구를 동일한 평가 날짜와 baseline으로 비교")
    setup = pd.DataFrame(
        [
            ["학습 범위", "구별 train pool 200일"],
            ["평가 날짜", "고정 7개 날짜"],
            ["후보 수", "Top-K 12"],
            ["수요예측", "1시간 horizon"],
            ["BC/Rollback", "사용하지 않음"],
            ["대표 성능", "Best checkpoint, 안정성은 Final checkpoint"],
        ],
        columns=["항목", "설정"],
    )
    add_table(slide, setup, 0.8, 1.35, 5.3, 4.25, font_size=12)
    add_card(slide, "평가 기준", "고정 7개 날짜의 평균 reward를 사용하고, 구별 MostImbalanced baseline 대비 Delta를 계산한다.", 6.6, 1.45, 5.85, 1.4, accent=BLUE)
    add_card(slide, "대표 성능", "Best checkpoint는 도달 가능한 성능, Final checkpoint는 학습 후반 안정성을 의미한다.", 6.6, 3.05, 5.85, 1.4, accent=GREEN)
    add_card(slide, "주의점", "구별 reward scale 차이가 크므로 절대 Delta와 Delta%를 함께 해석한다.", 6.6, 4.65, 5.85, 1.4, accent=AMBER)

    # 12. Hyperparameters
    slide = blank(prs)
    add_title(slide, "11. 주요 Hyperparameter", "알고리즘별 핵심 설정을 명시해 재현성을 높임")
    hp = pd.DataFrame(
        [
            ["REINFORCE", "γ=0.99, hidden=256, lrπ=3e-4, lrV=1e-3"],
            ["A2C", "γ=0.99, hidden=256, lrπ=1e-4, lrV=3e-4, batch=32"],
            ["PPO", "lr=1e-4, clip=0.1, ent=0.003, target_kl=0.03, n_steps=256"],
            ["Double DQN", "Double=True, Dueling=True, reward_scale=0.01, eps 0.3→0.02"],
            ["LinUCB", "α=0.5, l2=1.0, reward_scale=0.01"],
        ],
        columns=["Algorithm", "주요 설정"],
    )
    add_table(slide, hp, 0.6, 1.25, 12.15, 3.15, font_size=10)
    add_card(slide, "구조 하이퍼파라미터", "Top-K 후보 수는 단순 출력 옵션이 아니라 action space 크기를 바꾸는 핵심 설정이다.", 0.9, 4.75, 5.75, 1.15, accent=BLUE)
    add_card(slide, "실험 해석", "PPO와 DQN은 reward scale, clip range, target KL, Top-K 크기에 민감했다.", 6.95, 4.75, 5.45, 1.15, accent=AMBER)

    # 13. Ablation and variation
    slide = blank(prs)
    add_title(slide, "12. Hyperparameter / Ablation 실험", "Top-K 크기와 VAE latent feature를 추가로 비교")
    baseline_gangnam = -531.69
    dqn_rows = []
    for name, k, path in [
        ("DQN probe", 3, "logs/dqn_dqn_topk3_probe_dqn_강남구/history.npy"),
        ("DQN probe", 5, "logs/dqn_dqn_topk5_probe_dqn_강남구/history.npy"),
        ("DQN full", 12, "logs/dqn_interactive_dqn_강남구/history.npy"),
    ]:
        best, final, point = history_best_final(ROOT / path)
        if best is not None:
            dqn_rows.append([name, f"Top-K {k}", fmt_num(best), fmt_num(best - baseline_gangnam), f"{int(point):,} step"])
    dqn_ablation = pd.DataFrame(dqn_rows, columns=["실험", "설정", "Best", "Best Δ", "Best point"])
    add_table(slide, dqn_ablation, 0.65, 1.35, 6.0, 2.2, font_size=10)

    vae_summary = pd.DataFrame(
        [
            ["VAE-REINFORCE", "25개 구", int((vae["best_delta"] > 0).sum()), fmt_num(vae["best_delta"].mean()), fmt_num(vae["final_delta"].mean())],
            ["기본 REINFORCE", "25개 구", int(algo.loc[algo["algorithm"] == "REINFORCE", "best_win_districts"].iloc[0]), fmt_num(algo.loc[algo["algorithm"] == "REINFORCE", "mean_best_delta"].iloc[0]), fmt_num(algo.loc[algo["algorithm"] == "REINFORCE", "mean_final_delta"].iloc[0])],
        ],
        columns=["실험", "범위", "Best 승리 구", "Mean Best Δ", "Mean Final Δ"],
    )
    add_table(slide, vae_summary, 6.95, 1.35, 5.65, 2.2, font_size=10)
    add_card(slide, "해석", "Top-K를 줄이면 DQN의 강남구 결과는 개선됐지만, 너무 작은 K가 항상 좋은 것은 아니다. VAE는 일부 구에서는 도움이 되었지만 평균적으로는 안정적 개선을 만들지 못했다.", 0.8, 4.3, 11.8, 1.35, accent=GREEN)

    # 14. Seed and confidence interval
    slide = blank(prs)
    add_title(slide, "13. Seed / 신뢰구간에 대한 처리", "평가 기준 중 seed 반복 항목은 한계로 명시")
    add_card(slide, "현재 수행한 것", "서울 25개 구 전체를 같은 7개 평가일로 반복 평가하고, 구별 분포와 IQR을 학습곡선에 표시했다.", 0.85, 1.45, 5.6, 1.6, accent=GREEN)
    add_card(slide, "수행하지 못한 것", "동일 구·동일 설정에서 random seed만 바꾼 3회 이상 반복 실험과 confidence interval은 아직 포함하지 못했다.", 6.9, 1.45, 5.6, 1.6, accent=RED)
    add_card(slide, "보고서 해석 원칙", "따라서 본 결과는 seed confidence interval이 아니라 25개 구에 대한 cross-district robustness로 해석한다.", 0.85, 3.65, 11.65, 1.35, accent=BLUE)
    add_text(slide, "후속 보강: 대표 3개 구 × seed 3개 × A2C/PPO/DQN으로 평균±95% CI 작성", 1.2, 5.7, 10.9, 0.4, size=17, bold=True, color=TITLE, align=PP_ALIGN.CENTER)

    # 15. Overall results
    slide = blank(prs)
    add_title(slide, "14. 전체 결과 요약", "A2C가 Best와 Final 모두 가장 안정적인 평균 성능")
    out = algo.copy()
    out = out[["algorithm", "best_win_districts", "final_win_districts", "mean_best_delta", "mean_final_delta", "mean_best_delta_pct", "mean_final_delta_pct"]]
    out.columns = ["Algorithm", "Best 승리 구", "Final 승리 구", "Mean Best Δ", "Mean Final Δ", "Mean Best Δ%", "Mean Final Δ%"]
    out["Algorithm"] = out["Algorithm"].replace({"BANDIT": "LinUCB Bandit"})
    out["Mean Best Δ"] = out["Mean Best Δ"].map(lambda x: fmt_num(x))
    out["Mean Final Δ"] = out["Mean Final Δ"].map(lambda x: fmt_num(x))
    out["Mean Best Δ%"] = out["Mean Best Δ%"].map(lambda x: fmt_num(x))
    out["Mean Final Δ%"] = out["Mean Final Δ%"].map(lambda x: fmt_num(x))
    add_table(slide, out, 0.45, 1.25, 12.45, 3.1, font_size=9)
    add_card(slide, "핵심 결과", "A2C: 17/25개 구에서 Best 기준 baseline 초과\nREINFORCE/PPO: 가능성은 있으나 Final 안정성 낮음\nDQN/Bandit: 현재 구조에서 추가 개선 필요", 1.1, 4.85, 11.1, 1.3, accent=GREEN)

    # 16. Scorecard figure
    slide = blank(prs)
    add_title(slide, "15. 구별 성능 Scorecard")
    add_image(slide, FIG_DIR / "current_algorithm_delta_distribution.png", 0.45, 1.15, 12.4, 5.95)

    # 17. Learning curves
    slide = blank(prs)
    add_title(slide, "16. 학습곡선", "Best checkpoint만 보지 않고 학습 중 평가 reward 변화도 함께 확인")
    add_image(slide, FIG_DIR / "current_learning_curves.png", 0.65, 1.2, 12.0, 5.35)
    add_text(slide, "Best는 가능성, Final은 안정성이다. 두 값 차이가 크면 학습 후반 policy 유지가 어렵다는 의미다.", 1.15, 6.6, 11.0, 0.3, size=12, color=MUTED, align=PP_ALIGN.CENTER)

    # 18. Best/Worst
    slide = blank(prs)
    add_title(slide, "17. Best/Worst 구 분석")
    add_image(slide, FIG_DIR / "current_best_worst_learning_curves.png", 0.55, 1.25, 7.2, 5.45)
    add_image(slide, FIG_DIR / "current_best_worst_causal_scatter.png", 7.8, 1.35, 4.9, 4.85)
    add_text(slide, "큰 수요 규모, 낮은 forecast coverage, reward scale 차이가 성능 편차의 후보 원인이다.", 0.85, 6.65, 11.8, 0.28, size=11, color=MUTED, align=PP_ALIGN.CENTER)

    # 19. Algorithm interpretation
    slide = blank(prs)
    add_title(slide, "18. 알고리즘별 해석", "실험 결과를 RL 관점으로 정리")
    add_card(slide, "A2C", "TD update로 REINFORCE보다 빠르게 학습하며, 이번 실험에서 평균 성능과 Final 안정성이 가장 좋았다.", 0.75, 1.35, 3.75, 1.65, accent=GREEN)
    add_card(slide, "PPO", "clipped objective로 policy 변화 폭을 제한하지만, Top-K rank action이 매 step 바뀌어 안정성이 충분하지 않았다.", 4.8, 1.35, 3.75, 1.65, accent=AMBER)
    add_card(slide, "DQN", "Double DQN을 사용했지만 delayed reward와 rank action에서 Q-value target 학습이 어려웠다.", 8.85, 1.35, 3.75, 1.65, accent=RED)
    add_card(slide, "REINFORCE", "구별로 좋은 결과가 있었지만 Monte Carlo update 특성상 variance가 컸다.", 2.7, 3.55, 3.75, 1.65, accent=BLUE)
    add_card(slide, "LinUCB Bandit", "매우 빠르지만 다음 state 이후의 장기 재배치 효과를 보지 못한다.", 6.9, 3.55, 3.75, 1.65, accent=BLUE)

    # 20. Visualization demo
    slide = blank(prs)
    add_title(slide, "19. Replay Viewer", "학습된 정책이 episode 안에서 어떤 정류소를 선택하는지 시각화")
    add_card(slide, "시각화 대상", "정류소 위치, 트럭 위치, Top-K 후보, 후보별 선택 확률, 누적 reward, stockout/full을 함께 표시한다.", 0.8, 1.3, 5.1, 4.6, accent=BLUE)
    add_card(slide, "발표에서 보여줄 포인트", "1. policy가 후보 정류소 중 하나를 선택\n2. 트럭 이동에 따라 reward와 실패 지표 변화\n3. 표 결과가 실제 행동으로 어떻게 나타나는지 확인", 6.35, 1.3, 6.1, 4.6, accent=GREEN)

    # 21. Discussion
    slide = blank(prs)
    add_title(slide, "20. 토의", "이번 실험의 의미와 한계")
    add_bullets(
        slide,
        [
            "의미: 알고리즘만 바꾸는 것보다 state와 action 구조 설계가 성능에 더 직접적으로 영향을 주었다.",
            "한계: 25개 구 전체에서 Top-K 없음 ablation을 같은 비용으로 수행하지는 못했다.",
            "한계: VAE와 LinUCB Bandit은 탐색적 실험이며, seed 반복과 추가 hyperparameter ablation이 더 필요하다.",
            "향후: 실시간 재고 스냅샷, 다중 트럭, 구 간 이동, depot 위치를 포함하면 실제 운영 문제에 가까워진다.",
        ],
        0.9,
        1.35,
        11.7,
        4.5,
        size=17,
    )

    # 22. Conclusion
    slide = blank(prs)
    add_title(slide, "21. 결론", "따릉이 재배치 RL에서 핵심은 상태와 행동 구조 설계")
    add_card(slide, "1", "문제를 stockout/full/이동 비용 최소화로 정의하고 reward로 연결했다.", 0.8, 1.55, 3.7, 1.7, accent=BLUE)
    add_card(slide, "2", "1시간 수요예측과 capacity를 state에 추가해 미래 재고 위험을 보게 했다.", 4.8, 1.55, 3.7, 1.7, accent=GREEN)
    add_card(slide, "3", "Top-K 후보 action으로 탐색 공간을 줄이고 여러 RL 알고리즘을 비교했다.", 8.8, 1.55, 3.7, 1.7, accent=AMBER)
    add_text(slide, "현재 protocol에서는 A2C가 가장 안정적이었고, DQN/Bandit은 구조적·하이퍼파라미터 개선 여지가 컸다.", 1.0, 4.65, 11.4, 0.65, size=21, bold=True, color=TITLE, align=PP_ALIGN.CENTER)
    add_text(slide, "감사합니다", 5.1, 6.15, 3.1, 0.4, size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    ts = datetime.now().strftime("%Y-%m-%d_%H%M%S")
    out = OUT_DIR / f"ddareungi_rl_project_presentation_{ts}.pptx"
    prs.save(out)
    print(out)


if __name__ == "__main__":
    main()
